#!/usr/bin/env python3
"""Generate a Mac-friendly PPTX deck with standard layouts."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, item in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = item
        p.level = 0


def add_two_column_slide(
    prs: Presentation, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title

    left = slide.placeholders[1].text_frame
    left.clear()
    for idx, item in enumerate([left_title, *left_bullets]):
        p = left.paragraphs[0] if idx == 0 else left.add_paragraph()
        p.text = item
        p.level = 0 if idx == 0 else 1

    right = slide.placeholders[2].text_frame
    right.clear()
    for idx, item in enumerate([right_title, *right_bullets]):
        p = right.paragraphs[0] if idx == 0 else right.add_paragraph()
        p.text = item
        p.level = 0 if idx == 0 else 1


def build_deck() -> Presentation:
    prs = Presentation()

    add_title_slide(
        prs,
        "Comparatif planning - Orgue IAMS",
        "Ancien site vs nouveau prototype (version Mac-compatible)",
    )

    add_bullets_slide(
        prs,
        "1 - Architecture generale",
        [
            "Navigateur utilisateur -> application planning -> backend securise",
            "Application web: consultation, reservation, edition, administration",
            "Backend Supabase: auth, PostgreSQL, controle d'acces, fonctions serveur",
            "Integrations externes: Google Calendar et notifications e-mail",
        ],
    )

    add_two_column_slide(
        prs,
        "2 - Technologies utilisees",
        "Front",
        [
            "HTML / CSS / JavaScript (modules ES)",
            "FullCalendar",
            "Tailwind CSS + DaisyUI",
            "PWA",
        ],
        "Backend",
        [
            "Supabase Auth",
            "PostgreSQL",
            "Supabase Edge Functions",
            "Integrations Google Calendar + e-mail",
        ],
    )

    add_bullets_slide(
        prs,
        "3 - Hebergement actuel",
        [
            "Front: GitHub Pages (site statique HTTPS)",
            "Base de donnees: Supabase PostgreSQL",
            "Logique serveur: Supabase Edge Functions",
            "Cible a valider avec le SI (securite, RGPD, exploitation)",
        ],
    )

    add_two_column_slide(
        prs,
        "4 - Fonctionnel utilisateur",
        "Nouveau planning",
        [
            "Edition des creneaux fluide",
            "Experience mobile/tablette",
            "Types de creneaux clairs",
            "Consultation Google Calendar",
        ],
        "Ancien planning",
        [
            "Ergonomie web historique",
            "Fonctions degradees avec le temps",
            "Peu d'integration externe",
            "Maintenance difficile",
        ],
    )

    add_two_column_slide(
        prs,
        "5 - Pedagogie",
        "Nouveau planning",
        [
            "Type Cours distinct",
            "Association cours <-> eleves",
            "Gestion semaines A/B",
            "Lisibilite pedagogique amelioree",
        ],
        "Ancien planning",
        [
            "Pas de structure cours/eleves",
            "Intitules parfois ambigus",
            "Pas de logique A/B",
            "Risque de confusion",
        ],
    )

    add_bullets_slide(
        prs,
        "6 - Equite de reservation",
        [
            "Quota eleve travail perso: 2h/semaine (parametrable)",
            "Fenetre max reservation eleve: 15 jours (parametrable)",
            "Anti-chevauchement des reservations",
            "Controle des suppressions pour limiter les abus",
        ],
    )

    add_bullets_slide(
        prs,
        "7 - Securite",
        [
            "HTTPS maintenu",
            "Authentification modernisee (session/token)",
            "Droits renforces cote serveur et base de donnees",
            "Administration des comptes encadree (roles, suspension, mots de passe)",
            "Revue SI a finaliser pour policy mots de passe et RGPD detaille",
        ],
    )

    add_bullets_slide(
        prs,
        "8 - Statistiques deja disponibles",
        [
            "Occupation globale (cours, travail perso, concert, total)",
            "Indicateurs par eleve (creneaux, heures, moyenne)",
            "Graphique journalier des heures de travail perso",
            "Extensions possibles selon besoins de pilotage",
        ],
    )

    add_bullets_slide(
        prs,
        "9 - Comparatif synthetique (etoiles)",
        [
            "Authentification: Ancien 2/5 | Nouveau 4/5",
            "Ergonomie creneaux: Ancien 2/5 | Nouveau 5/5",
            "Mobile/tablette: Ancien 1/5 | Nouveau 4/5",
            "Pedagogie (cours/eleves/A-B): Ancien 1/5 | Nouveau 5/5",
            "Maintenance/gouvernance: Ancien 1/5 | Nouveau 5/5",
            "Securite globale: Ancien 2/5 | Nouveau 4/5",
        ],
    )

    add_bullets_slide(
        prs,
        "10 - Conclusion",
        [
            "Le nouveau planning est operationnel, plus lisible et plus pilotable",
            "Le socle est plus robuste et plus securise que l'ancien",
            "Prochaine etape: cadrage SI (securite/RGPD/exploitation)",
        ],
    )

    return prs


def main() -> None:
    target = (
        Path(__file__).resolve().parents[1]
        / "docs"
        / "presentation-directrice-comparatif-planning-mac.pptx"
    )
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(target)
    print(f"PPTX generated: {target}")


if __name__ == "__main__":
    main()
