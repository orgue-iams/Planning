#!/usr/bin/env python3
"""Generate the director presentation as a PPTX file."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

COLOR_BG = RGBColor(255, 255, 255)
COLOR_TITLE = RGBColor(15, 23, 42)
COLOR_KICKER = RGBColor(37, 99, 235)
COLOR_BODY = RGBColor(30, 41, 59)
COLOR_MUTED = RGBColor(71, 85, 105)
COLOR_BOX_BG = RGBColor(251, 253, 255)
COLOR_BOX_BORDER = RGBColor(219, 225, 234)

FONT_MAIN = "Segoe UI"


def add_kicker(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.8), Inches(0.45))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_KICKER
    p.font.name = FONT_MAIN


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.8), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE
    p.font.name = FONT_MAIN


def add_subtitle(slide, text: str, top: float = 1.8) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.8), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(19)
    p.font.color.rgb = COLOR_BODY
    p.font.name = FONT_MAIN


def add_bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 22,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_BODY
        p.font.name = FONT_MAIN
        p.space_after = Pt(7)


def add_text_block(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_BODY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_MAIN
    p.alignment = align


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    lines: list[str],
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BOX_BG
    shape.line.color.rgb = COLOR_BOX_BORDER
    shape.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.12), Inches(width - 0.4), Inches(0.35)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TITLE
    title_p.font.name = FONT_MAIN

    content_box = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.52), Inches(width - 0.4), Inches(height - 0.62)
    )
    content_tf = content_box.text_frame
    content_tf.clear()
    content_tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = content_tf.paragraphs[0] if idx == 0 else content_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BODY
        p.font.name = FONT_MAIN
        p.space_after = Pt(4)


def add_footer(slide, page: int, total: int) -> None:
    add_text_block(
        slide,
        "Presentation Direction - Planning Orgue IAMS",
        left=0.8,
        top=7.06,
        width=6.5,
        height=0.3,
        font_size=11,
        color=COLOR_MUTED,
    )
    add_text_block(
        slide,
        f"{page} / {total}",
        left=11.7,
        top=7.06,
        width=1.1,
        height=0.3,
        font_size=11,
        color=COLOR_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG
    return slide


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total = 11

    # 1
    slide = new_slide(prs)
    add_kicker(slide, "Orgue IAMS - Presentation Direction")
    add_title(slide, "Comparatif planning")
    add_subtitle(slide, "Ancien site vs Nouveau prototype", top=1.55)
    add_text_block(
        slide,
        "Objectif : montrer des resultats deja actifs, rassurer sur la securite, "
        "puis preparer une revue SI constructive.",
        left=0.8,
        top=2.4,
        width=11.8,
        height=0.9,
        font_size=20,
    )
    add_card(
        slide,
        left=0.8,
        top=3.3,
        width=11.8,
        height=2.1,
        title="Contexte",
        lines=[
            "Le nouveau planning conserve l'essentiel et corrige les fragilites "
            "de l'ancien site.",
            "Le prototype est deja fonctionnel et reste adaptable selon les retours Direction/SI.",
        ],
    )
    add_footer(slide, 1, total)

    # 2
    slide = new_slide(prs)
    add_kicker(slide, "1 - Vue d'ensemble")
    add_title(slide, "Architecture generale du planning")
    add_card(
        slide,
        left=0.8,
        top=1.8,
        width=11.8,
        height=1.0,
        title="Schema simple",
        lines=["navigateur utilisateur -> application planning -> services backend securises"],
    )
    add_card(
        slide,
        left=0.8,
        top=3.0,
        width=3.75,
        height=2.7,
        title="Application web",
        lines=["Consultation", "Reservation", "Edition", "Administration"],
    )
    add_card(
        slide,
        left=4.79,
        top=3.0,
        width=3.75,
        height=2.7,
        title="Backend Supabase",
        lines=["Authentification", "Base PostgreSQL", "Controle d'acces", "Fonctions serveur"],
    )
    add_card(
        slide,
        left=8.78,
        top=3.0,
        width=3.82,
        height=2.7,
        title="Integrations externes",
        lines=["Google Calendar", "Notifications e-mail"],
    )
    add_footer(slide, 2, total)

    # 3
    slide = new_slide(prs)
    add_kicker(slide, "2 - Technologies")
    add_title(slide, "Technos mises en oeuvre (prototype actuel)")
    add_card(
        slide,
        left=0.8,
        top=1.9,
        width=5.7,
        height=4.8,
        title="Front",
        lines=[
            "HTML / CSS / JavaScript (modules ES)",
            "FullCalendar (grille planning)",
            "Tailwind CSS + DaisyUI (interface)",
            "PWA (usage web fluide, mobile/tablette)",
        ],
    )
    add_card(
        slide,
        left=6.9,
        top=1.9,
        width=5.7,
        height=4.8,
        title="Backend",
        lines=[
            "Supabase Auth (sessions et roles)",
            "PostgreSQL (planning et profils)",
            "Edge Functions Supabase (logique serveur)",
            "Integrations : Google Calendar, e-mail",
        ],
    )
    add_footer(slide, 3, total)

    # 4
    slide = new_slide(prs)
    add_kicker(slide, "3 - Hebergement actuel")
    add_title(slide, "Ou tourne le prototype aujourd'hui")
    add_card(
        slide,
        left=0.8,
        top=2.0,
        width=11.8,
        height=1.4,
        title="Site web planning (front)",
        lines=["GitHub Pages - site statique public HTTPS, pratique pour demonstration."],
    )
    add_card(
        slide,
        left=0.8,
        top=3.65,
        width=11.8,
        height=1.4,
        title="Base de donnees",
        lines=["Supabase PostgreSQL - utilisateurs, creneaux, regles, statistiques."],
    )
    add_card(
        slide,
        left=0.8,
        top=5.3,
        width=11.8,
        height=1.2,
        title="Fonctions serveur",
        lines=[
            "Supabase Edge Functions - actions administratives, logique metier, integrations."
        ],
    )
    add_footer(slide, 4, total)

    # 5
    slide = new_slide(prs)
    add_kicker(slide, "4 - Fonctionnel utilisateur")
    add_title(slide, "Ce qui change au quotidien")
    add_card(
        slide,
        left=0.8,
        top=1.95,
        width=5.7,
        height=4.8,
        title="Nouveau planning",
        lines=[
            "Edition des creneaux fluide (deplacement/redimensionnement).",
            "Experience mobile/tablette moderne.",
            "Types clairs : Travail, Cours, Concert, Autre, Fermeture.",
            "Consultation Google Calendar (global + personnel).",
        ],
    )
    add_card(
        slide,
        left=6.9,
        top=1.95,
        width=5.7,
        height=4.8,
        title="Ancien planning",
        lines=[
            "Vue semaine surtout, ergonomie web classique.",
            "Fonctions historiques degradees (titres, mails, stats).",
            "Pas d'integration Google structuree.",
            "Maintenance devenue difficile.",
        ],
    )
    add_footer(slide, 5, total)

    # 6
    slide = new_slide(prs)
    add_kicker(slide, "5 - Pedagogie")
    add_title(slide, "Cours, eleves inscrits, semaines A/B")
    add_card(
        slide,
        left=0.8,
        top=1.95,
        width=5.7,
        height=4.8,
        title="Nouveau planning",
        lines=[
            "Type Cours distinct des reservations simples.",
            "Association cours <-> eleves lisible.",
            "Gestion des semaines A/B.",
            "Suivi plus fin des activites.",
        ],
    )
    add_card(
        slide,
        left=6.9,
        top=1.95,
        width=5.7,
        height=4.8,
        title="Ancien planning",
        lines=[
            "Pas de structure cours/eleves.",
            "Intitules parfois insuffisants.",
            "Pas de logique semaine A/B.",
            "Risque d'erreurs d'interpretation.",
        ],
    )
    add_footer(slide, 6, total)

    # 7
    slide = new_slide(prs)
    add_kicker(slide, "6 - Equite de reservation")
    add_title(slide, "Regles actives et parametrables")
    add_bullets(
        slide,
        [
            "Actif maintenant : quota eleve travail perso = 2h/semaine (modifiable).",
            "Actif maintenant : fenetre max de reservation eleve = 15 jours (modifiable).",
            "Actif maintenant : controle de suppression pour limiter les abus.",
            "Actif maintenant : anti-chevauchement des reservations.",
            "Avant, ces regles etaient surtout orales ; maintenant elles sont concretes et ajustables.",
        ],
        left=0.9,
        top=2.0,
        width=11.7,
        height=4.5,
        font_size=22,
    )
    add_footer(slide, 7, total)

    # 8
    slide = new_slide(prs)
    add_kicker(slide, "7 - Securite (message simple)")
    add_title(slide, "Pourquoi le nouveau socle est plus sur")
    add_bullets(
        slide,
        [
            "HTTPS maintenu (comme avant).",
            "Authentification modernisee : session/token, moins de logique fragile.",
            "Droits renforces : controles cote serveur et cote base.",
            "Administration encadree : roles, suspension, gestion des mots de passe.",
            "Gouvernance : versionning GitHub et suivi des corrections/evolutions.",
            "A cadrer avec le SI : politique mots de passe, sessions, RGPD detaille.",
        ],
        left=0.9,
        top=2.0,
        width=11.7,
        height=4.7,
        font_size=20,
    )
    add_footer(slide, 8, total)

    # 9
    slide = new_slide(prs)
    add_kicker(slide, "8 - Statistiques")
    add_title(slide, "Indicateurs deja en place")
    add_bullets(
        slide,
        [
            "Occupation globale : Cours / Travail perso (eleves) / Concert / Total.",
            "Par eleve : nombre de creneaux, heures, moyenne heures/semaine.",
            "Graphique journalier des heures de travail perso.",
            "Evolutif : ajout d'autres indicateurs selon les besoins de pilotage.",
            "Ancien site : module stats historique devenu non fonctionnel.",
        ],
        left=0.9,
        top=2.0,
        width=11.7,
        height=4.7,
        font_size=22,
    )
    add_footer(slide, 9, total)

    # 10
    slide = new_slide(prs)
    add_kicker(slide, "9 - Comparatif synthetique")
    add_title(slide, "Tableau a etoiles (lecture direction)")
    add_bullets(
        slide,
        [
            "Authentification/mots de passe : Ancien 2/5 | Nouveau 4/5",
            "Gestion creneaux & ergonomie : Ancien 2/5 | Nouveau 5/5",
            "Mobile/tablette : Ancien 1/5 | Nouveau 4/5",
            "Cours + eleves + semaines A/B : Ancien 1/5 | Nouveau 5/5",
            "Visibilite Google (global + perso) : Ancien 1/5 | Nouveau 5/5",
            "Maintenance/gouvernance : Ancien 1/5 | Nouveau 5/5",
            "Securite globale : Ancien 2/5 | Nouveau 4/5",
        ],
        left=0.9,
        top=2.0,
        width=11.7,
        height=4.8,
        font_size=20,
    )
    add_footer(slide, 10, total)

    # 11
    slide = new_slide(prs)
    add_kicker(slide, "10 - Conclusion")
    add_title(slide, "Message propose pour validation de principe")
    add_text_block(
        slide,
        "Le nouveau planning est operationnel, plus lisible, plus pilotable "
        "et plus securise que l'ancien.",
        left=0.8,
        top=2.0,
        width=11.8,
        height=1.0,
        font_size=26,
        bold=True,
    )
    add_bullets(
        slide,
        [
            "Base serieuse pour la revue SI (securite, RGPD, exploitation).",
            "Decision attendue : accord de principe pour passage en revue SI.",
            "Etape suivante : cadrage technique/securite avec le SI.",
        ],
        left=0.9,
        top=3.3,
        width=11.7,
        height=2.7,
        font_size=22,
    )
    add_footer(slide, 11, total)

    return prs


def main() -> None:
    target = Path(__file__).resolve().parents[1] / "docs" / "presentation-directrice-comparatif-planning.pptx"
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(target)
    print(f"PPTX generated: {target}")


if __name__ == "__main__":
    main()
